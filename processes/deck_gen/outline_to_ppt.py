#!/usr/bin/env python3
"""
outline_to_ppt.py
-----------------
Render a PowerPoint from an indices-backed, placeholder-exact outline.json.

The outline.json must have:
{
  "deck_meta": {...},
  "deck_slides": [
    {
      "layout_name": "...",        # optional, informational
      "master_index": 0,           # required
      "layout_index": 13,          # required
      "placeholders": {            # required; keys must match placeholder 'name' in template_map.json
        "title_main": "Agenda",
        "body_bullets": "- One\n- Two\n- Three"
      }
    },
    ...
  ]
}

Usage:
  python outline_to_ppt.py \
    --template templates/your_template.pptx \
    --template-map templates/template_map.json \
    --outline output/outline.json \
    --out output/generated_deck.pptx
"""
from __future__ import annotations

import argparse
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Tuple, Optional

from pptx import Presentation
from pptx.util import Pt

# pptx placeholder type ids that accept text
TEXT_TYPES = {1, 2, 3, 4, 7}
PICTURE_TYPE = 18


def load_json(p: Path) -> Any:
    with p.open("r", encoding="utf-8") as f:
        return json.load(f)


def build_layout_maps(template_map: Dict[str, Any]) -> Dict[int, Dict[str, Any]]:
    by_idx: Dict[int, Dict[str, Any]] = {}
    for entry in template_map.get("layouts", []):
        li = int(entry.get("layout_index", entry.get("index", -1)))
        if li < 0:
            continue
        phs = entry.get("raw_placeholders") or []
        name_to_idx: Dict[str, int] = {}
        idx_to_type: Dict[int, int] = {}
        for ph in phs:
            nm = (ph.get("name") or "").strip()
            if nm:
                name_to_idx[nm] = int(ph["idx"])
            if ph.get("idx") is not None and ph.get("type_id") is not None:
                idx_to_type[int(ph["idx"])] = int(ph["type_id"])
        by_idx[li] = {
            "name": entry.get("layout_name", f"layout_{li}"),
            "ph_name_to_idx": name_to_idx,
            "ph_idx_to_type": idx_to_type,
        }
    return by_idx


def find_placeholder_shape(slide, idx: int):
    for shp in slide.shapes:
        if not getattr(shp, "is_placeholder", False):
            continue
        try:
            if shp.placeholder_format.idx == idx:
                return shp
        except Exception:
            continue
    return None


def normalize_to_bullets(value: Any) -> Optional[List[str]]:
    if isinstance(value, list):
        return [str(v) for v in value]
    if not isinstance(value, str):
        return None
    s = value.strip()
    if not s:
        return []
    if "\n" in s or s.lstrip().startswith("-") or s.lstrip().startswith("•"):
        lines = [ln.strip() for ln in s.splitlines() if ln.strip()]
        cleaned = []
        for ln in lines:
            if ln.startswith("- "):
                cleaned.append(ln[2:].strip())
            elif ln.startswith("• "):
                cleaned.append(ln[2:].strip())
            elif ln.startswith("-"):
                cleaned.append(ln[1:].strip())
            elif ln.startswith("•"):
                cleaned.append(ln[1:].strip())
            else:
                cleaned.append(ln)
        return cleaned
    return None


def set_text(shape, value: Any) -> None:
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    tf.clear()
    bullets = normalize_to_bullets(value)
    if bullets is not None:
        if not bullets:
            if len(tf.paragraphs) == 0:
                tf.add_paragraph()
            tf.paragraphs[0].text = ""
            return
        for i, item in enumerate(bullets):
            if i == 0 and len(tf.paragraphs) > 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(item)
            p.level = 0
            p.line_spacing = 1.1
            p.space_after = Pt(6)
        return
    if len(tf.paragraphs) == 0:
        tf.add_paragraph()
    p = tf.paragraphs[0]
    p.text = "" if value is None else str(value)
    p.level = 0
    p.line_spacing = 1.1
    p.space_after = Pt(0)


def insert_picture(shape, image_path: str) -> None:
    if not image_path or not os.path.exists(image_path):
        return
    try:
        shape.placeholder_format.insert_picture(image_path)
    except Exception:
        slide = shape.part.slide
        slide.shapes.add_picture(image_path, shape.left, shape.top,
                                 width=shape.width, height=shape.height)


def render(template_pptx: Path, template_map_path: Path,
           outline_path: Path, out_path: Path) -> None:
    template_map = load_json(template_map_path)
    outline = load_json(outline_path)

    layouts_by_index = build_layout_maps(template_map)
    prs = Presentation(str(template_pptx))

    for i, slide_spec in enumerate(outline.get("deck_slides", [])):
        li = slide_spec.get("layout_index")
        if li is None or int(li) not in layouts_by_index:
            continue
        li = int(li)
        slide_layout = prs.slide_layouts[li]
        slide = prs.slides.add_slide(slide_layout)

        ph_name_to_idx = layouts_by_index[li]["ph_name_to_idx"]
        ph_idx_to_type = layouts_by_index[li]["ph_idx_to_type"]

        content: Dict[str, Any] = slide_spec.get("placeholders", {})
        ci_map = {k.lower(): v for k, v in ph_name_to_idx.items()}

        for ph_name, value in content.items():
            idx = ph_name_to_idx.get(ph_name) or ci_map.get(ph_name.lower())
            if idx is None:
                continue
            shp = find_placeholder_shape(slide, idx)
            if shp is None:
                continue
            phtype = ph_idx_to_type.get(idx)
            if phtype == PICTURE_TYPE:
                img_path = ""
                if isinstance(value, str):
                    img_path = value
                elif isinstance(value, dict):
                    img_path = (value.get("image_path") or value.get("path")
                                or value.get("image") or "")
                insert_picture(shp, img_path)
            elif phtype in TEXT_TYPES or phtype is None:
                set_text(shp, value)
            else:
                if isinstance(value, (str, list)):
                    set_text(shp, value)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved deck → {out_path}")


def main():
    ap = argparse.ArgumentParser(description="Render a deck from outline.json using template_map.json")
    ap.add_argument("--template", type=Path, default=Path("templates/mx_powerpoint_template_v3.pptx"))
    ap.add_argument("--template-map", type=Path, default=Path("templates/template_map.json"))
    ap.add_argument("--outline", type=Path, default=Path("output/outline.json"))
    ap.add_argument("--out", type=Path, default=Path("output/generated_deck.pptx"))
    args = ap.parse_args()
    render(args.template, args.template_map, args.outline, args.out)


if __name__ == "__main__":
    main()
